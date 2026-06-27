"""Create the PPTX status example template."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
TEMPLATE = ROOT / "templates" / "pptx-status.pptx"
DEFAULT_TEXT_COLOR = RGBColor(15, 23, 42)


def main() -> None:
    TEMPLATE.parent.mkdir(parents=True, exist_ok=True)

    presentation = Presentation()
    _remove_default_slide(presentation)
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)

    _add_title_slide(presentation)
    _add_marker_slide(presentation, "{% for milestone in milestones %}")
    _add_milestone_slide(presentation)
    _add_marker_slide(presentation, "{% endfor %}")
    _add_next_step_slide(presentation)

    presentation.save(TEMPLATE)


def _remove_default_slide(presentation: Presentation) -> None:
    for slide_id in list(presentation.slides._sldIdLst):
        presentation.slides._sldIdLst.remove(slide_id)


def _blank_slide(presentation: Presentation):
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    background = slide.background.fill
    background.solid()
    background.fore_color.rgb = RGBColor(248, 250, 252)
    return slide


def _add_title_slide(presentation: Presentation) -> None:
    slide = _blank_slide(presentation)
    _add_kicker(slide, "Stencil example")
    _add_text(slide, "{{ project.name }}", 0.75, 1.0, 11.0, 0.8, 38, bold=True)
    _add_text(slide, "Owner: {{ project.owner }}", 0.8, 2.2, 11.0, 0.4, 20)
    _add_text(slide, "Status: {{ project.status }}", 0.8, 2.85, 11.0, 0.4, 20)
    _add_rule(slide, 0.8, 4.0, 4.0)
    _add_text(
        slide,
        "This is a template deck; render it to create a presentation.",
        0.8,
        4.35,
        10.8,
        0.4,
        17,
    )


def _add_marker_slide(presentation: Presentation, tag: str) -> None:
    slide = _blank_slide(presentation)
    _add_text(slide, tag, 0.8, 3.15, 11.5, 0.7, 24, color=RGBColor(71, 85, 105))


def _add_milestone_slide(presentation: Presentation) -> None:
    slide = _blank_slide(presentation)
    _add_kicker(slide, "Milestone")
    _add_text(slide, "{{ milestone.name }}", 0.75, 1.0, 11.0, 0.75, 34, bold=True)
    _add_status_pill(slide, "{{ milestone.status }}")
    _add_text(slide, "{{ milestone.summary }}", 0.85, 3.0, 10.9, 1.0, 23)


def _add_next_step_slide(presentation: Presentation) -> None:
    slide = _blank_slide(presentation)
    _add_kicker(slide, "Next step")
    _add_text(slide, "Open the rendered deck", 0.75, 1.25, 11.0, 0.75, 36, bold=True)
    _add_text(
        slide,
        "Confirm the milestone detail slide was repeated once for each item in the JSON data.",
        0.85,
        2.45,
        10.8,
        0.9,
        22,
    )


def _add_kicker(slide, text: str) -> None:
    _add_text(slide, text.upper(), 0.8, 0.45, 11.0, 0.3, 12, color=RGBColor(15, 118, 110))


def _add_rule(slide, left: float, top: float, width: float) -> None:
    shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(0.04))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(15, 118, 110)
    shape.line.fill.background()


def _add_status_pill(slide, text: str) -> None:
    shape = slide.shapes.add_shape(5, Inches(0.85), Inches(2.05), Inches(3.4), Inches(0.48))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(220, 252, 231)
    shape.line.color.rgb = RGBColor(22, 163, 74)
    frame = shape.text_frame
    frame.clear()
    paragraph = frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    run = paragraph.add_run()
    run.text = text
    run.font.size = Pt(15)
    run.font.bold = True
    run.font.color.rgb = RGBColor(21, 128, 61)


def _add_text(
    slide,
    text: str,
    left: float,
    top: float,
    width: float,
    height: float,
    size: int,
    *,
    bold: bool = False,
    color: RGBColor = DEFAULT_TEXT_COLOR,
) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.clear()
    paragraph = frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


if __name__ == "__main__":
    main()
